"""Documentos generados por MAX — Step 12 expandido.

GET  /docs                  — lista de documentos
POST /docs/generate         — genera documento con LLM (todos los formatos)
GET  /docs/{id}/download    — descarga el documento

Formatos soportados:
  Texto:   md, html, txt, csv
  Código:  py, java, cpp, js, ts, sql, css, json, yaml
  Office:  docx, xlsx, pptx
  Especial: excalidraw, canvas
"""

from __future__ import annotations

import io
import json
import textwrap
import uuid
from pathlib import Path
from typing import Any

from fastapi import APIRouter, HTTPException, status
from fastapi.responses import FileResponse
from pydantic import BaseModel

from app.chat.dispatcher import stream_response
from app.core.config import get_settings
from app.core.deps import CurrentUser, SupabaseDep
from app.kimi.client import generate_text as kimi_generate

router = APIRouter(prefix="/docs", tags=["docs"])

DOCS_DIR = Path("/tmp/max-docs")
DOCS_DIR.mkdir(exist_ok=True)

# ── Tipos de formato ──────────────────────────────────────────────────────────

TEXT_FORMATS: set[str] = {
    "md", "html", "txt", "csv",
    "py", "java", "cpp", "js", "ts", "sql", "css", "json", "yaml",
}
BINARY_FORMATS: set[str] = {"docx", "xlsx", "pptx"}
JSON_FORMATS: set[str] = {"excalidraw", "canvas"}

ALL_FORMATS = TEXT_FORMATS | BINARY_FORMATS | JSON_FORMATS

MIME_TYPES: dict[str, str] = {
    # texto
    "md": "text/markdown",
    "html": "text/html",
    "txt": "text/plain",
    "csv": "text/csv",
    # código
    "py": "text/x-python",
    "java": "text/x-java-source",
    "cpp": "text/x-c++src",
    "js": "application/javascript",
    "ts": "application/typescript",
    "sql": "application/sql",
    "css": "text/css",
    "json": "application/json",
    "yaml": "application/yaml",
    # office
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    # especiales
    "excalidraw": "application/json",
    "canvas": "application/json",
}


# ── Prompts por formato ───────────────────────────────────────────────────────

def _system_hint(fmt: str) -> str:
    """Instrucción de formato para el LLM."""
    hints: dict[str, str] = {
        "md": "Responde en Markdown limpio con headings, listas y bloques de código donde aplique.",
        "html": "Responde solo con el contenido del <body> en HTML semántico, sin <html>/<head>.",
        "txt": "Responde en texto plano bien formateado, sin markdown.",
        "csv": (
            "Responde SOLO con datos CSV válidos (primera fila = cabeceras)."
            " Sin explicaciones, solo CSV."
        ),
        "py": "Responde SOLO con código Python válido, comentado. Sin explicaciones fuera del código.",
        "java": "Responde SOLO con código Java válido y completo, con comentarios JavaDoc.",
        "cpp": "Responde SOLO con código C++ válido con includes necesarios y comentarios.",
        "js": "Responde SOLO con código JavaScript moderno (ES2022+), comentado.",
        "ts": "Responde SOLO con código TypeScript con tipos explícitos, comentado.",
        "sql": "Responde SOLO con SQL válido (PostgreSQL por defecto), con comentarios.",
        "css": "Responde SOLO con CSS válido, organizado por secciones comentadas.",
        "json": "Responde SOLO con JSON válido, bien estructurado. Sin markdown ni explicaciones.",
        "yaml": "Responde SOLO con YAML válido. Sin explicaciones.",
        "docx": (
            "Responde en Markdown bien estructurado con # para título, ## para secciones,"
            " ### para subsecciones, y listas con - para bullets."
        ),
        "xlsx": (
            "Responde SOLO con datos CSV (primera fila = cabeceras en español)."
            " Separa por comas. Sin explicaciones, solo los datos."
        ),
        "pptx": (
            "Estructura la respuesta como diapositivas. Usa ## para cada slide,"
            " su contenido como bullets con - . Empieza con ## Portada."
            " Máximo 6 slides, máximo 5 bullets por slide."
        ),
        "excalidraw": (
            "Responde en Markdown con # para título, ## para secciones y"
            " texto descriptivo. Será convertido a diagrama visual."
        ),
        "canvas": (
            "Responde en Markdown con # para el nodo principal, ## para nodos secundarios."
            " Cada sección se convertirá en un nodo del Canvas."
        ),
    }
    return hints.get(fmt, "Responde en el formato adecuado para el tipo de archivo.")


# ── Generadores binarios ──────────────────────────────────────────────────────

def _make_docx(title: str, content: str) -> bytes:
    from docx import Document  # lazy import
    from docx.shared import Pt

    doc = Document()
    doc.add_heading(title, 0)
    for line in content.split("\n"):
        stripped = line.rstrip()
        if stripped.startswith("### "):
            doc.add_heading(stripped[4:], 3)
        elif stripped.startswith("## "):
            doc.add_heading(stripped[3:], 2)
        elif stripped.startswith("# "):
            doc.add_heading(stripped[2:], 1)
        elif stripped.startswith(("- ", "* ")):
            doc.add_paragraph(stripped[2:], style="List Bullet")
        elif stripped:
            doc.add_paragraph(stripped)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_xlsx(title: str, content: str) -> bytes:
    import csv as csv_mod
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill

    wb = Workbook()
    ws = wb.active
    ws.title = title[:31]  # máx 31 chars en Excel

    reader = csv_mod.reader(io.StringIO(content.strip()))
    for row_idx, row in enumerate(reader, 1):
        for col_idx, value in enumerate(row, 1):
            cell = ws.cell(row=row_idx, column=col_idx, value=value.strip())
            if row_idx == 1:
                cell.font = Font(bold=True)
                cell.fill = PatternFill("solid", fgColor="4F46E5")
                cell.font = Font(bold=True, color="FFFFFF")

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _make_pptx(title: str, content: str) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation()

    # Slide de título
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if slide.placeholders[1:]:
        slide.placeholders[1].text = "Generado por MAX Studio"

    # Slides de contenido: parsear secciones ## …
    sections = content.split("\n## ")
    for section in sections:
        lines = [ln.rstrip() for ln in section.split("\n") if ln.strip()]
        if not lines:
            continue
        slide_title = lines[0].lstrip("# ").strip()
        bullets = [ln[2:].strip() for ln in lines[1:] if ln.startswith("- ")]
        body_lines = [ln for ln in lines[1:] if not ln.startswith("-")]

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_title

        tf = slide.placeholders[1].text_frame
        tf.word_wrap = True
        tf.text = "\n".join(bullets) if bullets else "\n".join(body_lines)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_excalidraw(title: str, content: str) -> str:
    elements = []
    y = 0
    lines = content.split("\n")
    for line in lines:
        stripped = line.rstrip()
        if not stripped:
            y += 20
            continue
        font_size = 24 if stripped.startswith("# ") else 18 if stripped.startswith("## ") else 14
        text = stripped.lstrip("# -* ").strip()
        elements.append({
            "id": uuid.uuid4().hex[:8],
            "type": "text",
            "x": 40 if not stripped.startswith("#") else 20,
            "y": y,
            "width": 720,
            "height": font_size + 8,
            "angle": 0,
            "strokeColor": "#1e1e2e",
            "backgroundColor": "transparent",
            "fillStyle": "solid",
            "strokeWidth": 2,
            "roughness": 1,
            "opacity": 100,
            "text": text,
            "fontSize": font_size,
            "fontFamily": 1,
            "textAlign": "left",
            "verticalAlign": "top",
            "version": 1,
        })
        y += font_size + 12
    return json.dumps({
        "type": "excalidraw",
        "version": 2,
        "source": "MAX Studio",
        "elements": elements,
        "appState": {"gridSize": None, "viewBackgroundColor": "#ffffff"},
        "files": {},
    }, indent=2)


def _make_canvas(title: str, content: str) -> str:
    nodes = []
    x, y = 0, 0
    sections = content.split("\n## ")

    # Nodo principal
    main_text = sections[0].strip()
    nodes.append({
        "id": "main",
        "type": "text",
        "text": f"# {title}\n\n{main_text}",
        "x": x, "y": y,
        "width": 500,
        "height": max(150, len(main_text.split("\n")) * 25),
    })

    # Nodos secundarios
    for i, section in enumerate(sections[1:], 1):
        lines = section.strip().split("\n", 1)
        sec_title = lines[0].strip()
        sec_body = lines[1].strip() if len(lines) > 1 else ""
        sec_x = x + 600 * (i % 3)
        sec_y = y + 400 * (i // 3)
        nodes.append({
            "id": f"node{i}",
            "type": "text",
            "text": f"## {sec_title}\n\n{sec_body}",
            "x": sec_x, "y": sec_y,
            "width": 400,
            "height": max(120, len(sec_body.split("\n")) * 22),
        })

    return json.dumps({"nodes": nodes, "edges": []}, indent=2, ensure_ascii=False)


# ── Endpoints ─────────────────────────────────────────────────────────────────

class GenerateRequest(BaseModel):
    title: str
    format: str = "md"
    prompt: str


@router.get("")
async def list_docs(user: CurrentUser, sb: SupabaseDep) -> list[dict[str, Any]]:  # noqa: ARG001
    return await sb.select_many(
        "generated_docs",
        columns="id,filename,mime_type,size_bytes,created_at",
        order="created_at.desc",
        limit=100,
    )


@router.post("/generate", status_code=status.HTTP_201_CREATED)
async def generate_doc(
    body: GenerateRequest,
    user: CurrentUser,  # noqa: ARG001
    sb: SupabaseDep,
) -> dict[str, Any]:
    settings = get_settings()
    fmt = body.format.lower().strip(". ")

    if fmt not in ALL_FORMATS:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Formato no soportado: {fmt}. Disponibles: {sorted(ALL_FORMATS)}",
        )

    hint = _system_hint(fmt)
    full_prompt = (
        f"Genera un documento titulado: {body.title!r}.\n\n"
        f"Instrucción: {body.prompt}\n\n"
        f"Formato requerido: {hint}"
    )

    # Formatos binarios/complejos → Kimi K2.6 (mejor estructura, 256K contexto)
    # Resto → Dispatch (gpt-120, más rápido y gratis)
    use_kimi = fmt in BINARY_FORMATS and bool(settings.moonshot_api_key)

    raw_content = ""
    if use_kimi:
        try:
            raw_content = await kimi_generate(
                prompt=full_prompt,
                system="Eres un asistente experto en generación de documentos profesionales.",
            )
        except Exception as exc:
            # Fallback a Dispatch si Kimi falla
            import logging
            logging.getLogger(__name__).warning("Kimi fallback to Dispatch: %s", exc)
            use_kimi = False

    if not use_kimi:
        async for token in stream_response(full_prompt, [], model=settings.default_model):
            raw_content += token

    if not raw_content.strip():
        raise HTTPException(
            status_code=status.HTTP_502_BAD_GATEWAY,
            detail="El LLM no devolvió contenido",
        )

    # Generar archivo según tipo
    safe_title = "".join(c if c.isalnum() or c in "-_ " else "_" for c in body.title)[:40]
    filename = f"{safe_title.replace(' ', '_')}_{uuid.uuid4().hex[:6]}.{fmt}"
    filepath = DOCS_DIR / filename

    try:
        if fmt == "docx":
            filepath.write_bytes(_make_docx(body.title, raw_content))
        elif fmt == "xlsx":
            filepath.write_bytes(_make_xlsx(body.title, raw_content))
        elif fmt == "pptx":
            filepath.write_bytes(_make_pptx(body.title, raw_content))
        elif fmt == "excalidraw":
            filepath.write_text(_make_excalidraw(body.title, raw_content), encoding="utf-8")
        elif fmt == "canvas":
            filepath.write_text(_make_canvas(body.title, raw_content), encoding="utf-8")
        else:
            filepath.write_text(raw_content, encoding="utf-8")
    except ImportError as exc:
        raise HTTPException(
            status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
            detail=f"Librería no instalada para .{fmt}: {exc}",
        ) from exc

    size = filepath.stat().st_size

    row = await sb.insert(
        "generated_docs",
        {
            "filename": filename,
            "mime_type": MIME_TYPES[fmt],
            "storage_path": str(filepath),
            "size_bytes": size,
        },
    )

    return {
        "id": (row or {}).get("id"),
        "filename": filename,
        "size_bytes": size,
        "preview": textwrap.shorten(raw_content, width=300, placeholder="…"),
    }


@router.get("/{doc_id}/download")
async def download_doc(
    doc_id: str,
    user: CurrentUser,  # noqa: ARG001
    sb: SupabaseDep,
) -> FileResponse:
    row = await sb.select_one("generated_docs", {"id": doc_id})
    if not row:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Documento no encontrado")

    path = Path(row["storage_path"])
    if not path.exists():
        raise HTTPException(status_code=status.HTTP_410_GONE, detail="Archivo eliminado del servidor")

    return FileResponse(
        path=str(path),
        media_type=row["mime_type"],
        filename=row["filename"],
    )
